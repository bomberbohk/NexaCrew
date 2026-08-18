"""SkillScript BASIC — enterprise BASIC-style scripting language for custom
skills & plugins.

A skill whose instructions start with ``PROGRAM <name>`` is treated as a
SkillScript program instead of a plain prompt. Programs are executed in a
sandboxed interpreter (step / time / IO limits, no arbitrary Python) and the
text they PRINT/OUTPUT becomes the dynamic instructions injected into the
agents — so skills can compute, branch, loop, call AI models, fetch HTTP
data, read/write files and keep persistent memory.

STATEMENTS (case-insensitive, optional line numbers, ' or REM for comments)
---------------------------------------------------------------------------
  PROGRAM name / END PROGRAM        program envelope
  LET x = expr        |  x = expr   assignment (aliases: SET, VAR, DECLARE)
  CONST NAME = expr                 constant (alias: DEFINE)
  DIM a(n)                          array of n empty slots
  SWAP a, b                         exchange two variables
  INC x [, n] | DEC x [, n]         increment / decrement (default 1)
  PUSH list, expr | POP var = list  stack operations on a list variable
  PRINT expr[; expr…]               append to program output
                                    (aliases: OUTPUT, SAY, ECHO, WRITE,
                                     WRITELN, DISPLAY, SHOW, TELL)
  TITLE expr | HEADER expr          markdown "# …" / "## …" headings
  BULLET expr                       markdown "* …" bullet line
  HR ["char"]                       horizontal rule line
  CLS                               clear all output produced so far
  INPUT var [, "prompt"]            read a named input from the run context
  IF expr THEN [stmt]               inline or block form
  ELSEIF expr THEN | ELSE | ENDIF   (END IF also accepted)
  SELECT CASE expr / CASE val[, v2] / CASE ELSE / END SELECT
  FOR i = a TO b [STEP c] … NEXT    counted loop
  WHILE expr … WEND                 conditional loop
  DO … LOOP WHILE expr|LOOP UNTIL expr
  BREAK | EXIT FOR | EXIT WHILE | CONTINUE
  label:  |  GOTO label (alias JUMP)  |  GOSUB label … RETURN
  SUB name … END SUB  |  CALL name  (aliases: PERFORM, INVOKE)
  ASK var = expr [SYSTEM expr] [USING "codex"|"claude"|"api"]   AI model call
                                    (aliases: QUERY, LLM, GPT)
  SUMMARIZE var = text [WORDS n]    AI: concise summary
  TRANSLATE var = text TO "lang"    AI: translation
  SENTIMENT var = text              AI: positive / negative / neutral
  CLASSIFY var = text INTO "a,b,c"  AI: pick one category
  EXTRACT var = text FIELDS "f,g"   AI: extract fields as JSON
  REWRITE var = text [STYLE expr]   AI: rewrite text in a style
  EXPAND var = notes [WORDS n]      AI: expand notes into prose
  OUTLINE var = topic [DEPTH n]     AI: markdown outline
  PROOFREAD var = text              AI: grammar / spelling fix
  KEYWORDS var = text [TOP n]       AI: keyword list
  TITLEGEN var = text [STYLE expr]  AI: generate a title
  REPLYGEN var = message [TONE e]   AI: draft a reply
  CODEGEN var = task [LANG "py"]    AI: generate code
  SQLGEN var = request [DIALECT e]  AI: generate SQL
  IMAGE "file.png", "description"   AI image generation (→ IMAGE_RESULT)
  GENERATE PDF|DOCX|XLSX|PPTX|HTML|TXT|MD|CSV|JSON|XML path, content
                                    create a real file (→ LAST_FILE)
  CHART "f.svg", "bar"|"line"|"pie", labels, values [, title]
                                    render an SVG chart file (→ LAST_FILE)
  HTTP var = GET expr               HTTP GET → var  (alias: FETCH)
  HTTP var = POST expr BODY expr    HTTP POST (JSON body)
  FILE READ var = expr              read text file (inside data dirs)
  FILE WRITE expr, expr             write text file
  FILE APPEND expr, expr            append text file
  MEM SET expr, expr                persistent key/value store (survives runs)
  MEM GET var = expr                read persistent store ("" if missing)
  EMAIL TO expr SUBJECT expr BODY expr    send email via configured SMTP
  LOG expr                          write to the execution trace
                                    (aliases: DEBUG, TRACE, NOTE, INFO, WARN)
  SLEEP expr                        pause, max 10 s (aliases: WAIT, PAUSE, DELAY)
  RANDOMIZE [seed]                  seed the random generator
  ASSERT expr [, expr]              abort with message when false
                                    (aliases: REQUIRE, VERIFY, EXPECT, CHECK)
  ON ERROR GOTO label | ON ERROR RESUME NEXT | ON ERROR STOP
  THROW expr  (aliases: FAIL, PANIC) raise a script error
  NOP | BEEP                        do nothing (placeholders)
  END | STOP  (aliases: HALT, QUIT, FINISH, TERMINATE)   finish successfully
  ABORT expr                        finish with an error

BUILT-IN FUNCTIONS — 230+, usable in any expression
---------------------------------------------------------------------------
Strings:
  LEN UPPER LOWER TRIM LTRIM RTRIM LEFT RIGHT MID REPLACE SPLIT JOIN INSTR
  LASTINSTR CONTAINS STARTSWITH ENDSWITH REPEAT FORMAT CHR ASC CAPITALIZE
  TITLECASE SWAPCASE PADLEFT PADRIGHT CENTER ZFILL SQUEEZE REVERSESTR
  COUNTSTR TRUNCATE SLUG CAMELCASE SNAKECASE KEBABCASE PASCALCASE WRAP
  INDENT QUOTE UNQUOTE BETWEEN BEFORE AFTER REMOVE INSERTSTR SPACES TAB NL
  ISALPHA ISDIGIT ISALNUM ISUPPER ISLOWER ISSPACE LEVENSHTEIN SIMILARITY
  LINES WORDS
Math:
  ABS INT ROUND MIN MAX SQR CBRT POW EXP LN LOG10 LOG2 SIN COS TAN ASIN
  ACOS ATAN ATAN2 SINH COSH TANH FLOOR CEIL SIGN FRAC GCD LCM FACT COMB
  PERM CLAMP LERP MAPRANGE DEG RAD HYPOT STR VAL PARSEINT PARSEFLOAT
Statistics:
  MEAN MEDIAN MODE STDEV VARIANCE PERCENTILE ZSCORE AVERAGE
Random:
  RND RANDINT CHOICE SHUFFLE SAMPLE UUID RANDSTR
Date & time:
  NOW TODAY TIMESTAMP YEAR MONTH DAY WEEKDAY HOUR MINUTE SECOND DATEADD
  DATEDIFF DATEFORMAT DAYOFYEAR WEEKNUM ISLEAPYEAR DAYSINMONTH QUARTER
  MONTHNAME EPOCH2DATE DATE2EPOCH AGE
Lists:
  ARRAY APPEND SORT SORTDESC REVERSE COUNT SUM FIRST LAST NTH INDEXOF
  SLICE UNIQUE FLATTEN ZIP FILL SEQ TAKE DROP CHUNK INTERSECT UNIONL
  DIFFERENCE PRODUCT
JSON:
  JSONPARSE JSONGET JSONSTR ISJSON KEYS VALUES HASKEY
Encoding & crypto:
  HASH MD5 SHA1 SHA256 SHA512 CRC32 B64ENCODE B64DECODE HEXENCODE HEXDECODE
  ROT13 URLENCODE URLDECODE HTMLESCAPE HTMLUNESCAPE
Regex:
  REGEXMATCH REGEXEXTRACT REGEXEXTRACTALL REGEXREPLACE
Conversion & formatting:
  BIN OCT HEXN TOBOOL CURRENCY PERCENTF COMMA FIXED ROMAN UNROMAN
Units:
  C2F F2C KM2MI MI2KM KG2LB LB2KG M2FT FT2M L2GAL GAL2L
Finance & geometry:
  PMT FV PV COMPOUND AREACIRCLE AREARECT AREATRI PYTHAG
Colors:
  RGB2HEX HEX2RGB
Validation:
  ISEMAIL ISURL ISIP ISPHONE ISUUID ISDATE ISNUMBER ISEMPTY
Logic & misc:
  IIF COALESCE DEFAULTVAL TYPEOF

OPERATORS: + - * / \\ (int div) MOD ^ & (concat) = <> < <= > >= AND OR NOT

BUILT-IN VARIABLES: SKILL_NAME, TARGET, USER_INPUT, DATE$, TIME$, NODE$,
ERR, ERRMSG, HTTP_STATUS, TRUE, FALSE, NULL, PI — plus every key passed in
the run context.

Sandbox limits: 20 000 statements, 120 s wall clock, 512 KB output, HTTP
responses capped at 1 MB, files restricted to the platform data folders,
SLEEP capped at 10 s, AI calls capped at 5 per run.
"""
from __future__ import annotations

import base64
import datetime as dt
import hashlib
import json
import math
import random
import re
import time
import urllib.parse
import urllib.request
from pathlib import Path
from typing import Any, Callable, Dict, List, Optional

MAX_STEPS = 20_000
MAX_SECONDS = 120.0
MAX_OUTPUT = 512 * 1024
MAX_HTTP_BYTES = 1024 * 1024
MAX_SLEEP = 10.0
MAX_AI_CALLS = 5

DATA_DIR = Path(__file__).resolve().parent.parent / "data"
MEM_FILE = DATA_DIR / "skillscript_memory.json"


class ScriptError(Exception):
    def __init__(self, message: str, line: int = 0):
        super().__init__(message)
        self.line = line


def is_script(text: str) -> bool:
    """True when a skill's instructions are a SkillScript BASIC program."""
    for ln in (text or "").splitlines():
        s = ln.strip()
        if not s or s.startswith("'") or s.upper().startswith("REM"):
            continue
        s = re.sub(r"^\d+\s+", "", s)  # optional BASIC line numbers
        return s.upper().startswith("PROGRAM ")
    return False


# ---------------------------------------------------------------- memory
def _mem_load() -> dict:
    try:
        if MEM_FILE.is_file():
            return json.loads(MEM_FILE.read_text(encoding="utf-8"))
    except Exception:
        pass
    return {}


def _mem_save(mem: dict) -> None:
    MEM_FILE.parent.mkdir(parents=True, exist_ok=True)
    MEM_FILE.write_text(json.dumps(mem, ensure_ascii=False, indent=1), encoding="utf-8")


# ------------------------------------------------------------- functions
def _jsonget(obj: Any, path: str) -> Any:
    cur = obj
    for part in str(path).split("."):
        if isinstance(cur, list):
            cur = cur[int(part)]
        elif isinstance(cur, dict):
            cur = cur.get(part, "")
        else:
            return ""
    return cur


def _val(s: Any) -> float:
    try:
        f = float(str(s).strip() or 0)
        return int(f) if f.is_integer() else f
    except ValueError:
        return 0


def _funcs() -> Dict[str, Callable]:
    return {
        "LEN": lambda x: len(x), "UPPER": lambda s: str(s).upper(),
        "LOWER": lambda s: str(s).lower(), "TRIM": lambda s: str(s).strip(),
        "LEFT": lambda s, n: str(s)[: int(n)], "RIGHT": lambda s, n: str(s)[-int(n):] if int(n) else "",
        "MID": lambda s, i, n=None: str(s)[int(i) - 1: (int(i) - 1 + int(n)) if n is not None else None],
        "REPLACE": lambda s, a, b: str(s).replace(str(a), str(b)),
        "SPLIT": lambda s, sep=",": str(s).split(str(sep)),
        "JOIN": lambda lst, sep=",": str(sep).join(str(x) for x in lst),
        "INSTR": lambda s, sub: str(s).find(str(sub)) + 1,
        "CONTAINS": lambda s, sub: str(sub) in str(s),
        "STARTSWITH": lambda s, p: str(s).startswith(str(p)),
        "ENDSWITH": lambda s, p: str(s).endswith(str(p)),
        "REPEAT": lambda s, n: str(s) * int(n),
        "FORMAT": lambda fmt, *a: str(fmt).format(*a),
        "CHR": lambda n: chr(int(n)), "ASC": lambda s: ord(str(s)[0]),
        "STR": lambda x: str(x), "VAL": _val,
        "ABS": abs, "INT": lambda x: int(x), "ROUND": lambda x, d=0: round(x, int(d)),
        "MIN": min, "MAX": max, "SQR": math.sqrt, "POW": pow,
        "RND": random.random, "RANDINT": lambda a, b: random.randint(int(a), int(b)),
        "NOW": lambda: dt.datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
        "TODAY": lambda: dt.date.today().isoformat(),
        "TIMESTAMP": lambda: int(time.time()),
        "YEAR": lambda: dt.date.today().year, "MONTH": lambda: dt.date.today().month,
        "DAY": lambda: dt.date.today().day,
        "WEEKDAY": lambda: dt.date.today().strftime("%A"),
        "HOUR": lambda: dt.datetime.now().hour, "MINUTE": lambda: dt.datetime.now().minute,
        "JSONPARSE": lambda s: json.loads(s), "JSONGET": _jsonget,
        "JSONSTR": lambda o: json.dumps(o, ensure_ascii=False),
        "ARRAY": lambda *a: list(a),
        "APPEND": lambda lst, x: (lst.append(x) or lst),
        "SORT": lambda lst: sorted(lst), "REVERSE": lambda lst: list(reversed(lst)),
        "COUNT": lambda lst: len(lst), "SUM": lambda lst: sum(_val(x) for x in lst),
        "FIRST": lambda lst: lst[0] if lst else "", "LAST": lambda lst: lst[-1] if lst else "",
        "SLICE": lambda lst, a, b=None: lst[int(a): int(b) if b is not None else None],
        "UNIQUE": lambda lst: list(dict.fromkeys(lst)),
        "LINES": lambda s: str(s).splitlines(), "WORDS": lambda s: str(s).split(),
        "ISNUMBER": lambda x: isinstance(x, (int, float)) or str(x).replace(".", "", 1).lstrip("-").isdigit(),
        "ISEMPTY": lambda x: not x,
        "IIF": lambda c, a, b: a if c else b,
        "TYPEOF": lambda x: type(x).__name__,
        "HASH": lambda s: hashlib.sha256(str(s).encode()).hexdigest()[:16],
        "B64ENCODE": lambda s: base64.b64encode(str(s).encode()).decode(),
        "B64DECODE": lambda s: base64.b64decode(str(s)).decode("utf-8", "replace"),
        "URLENCODE": lambda s: urllib.parse.quote(str(s)),
        "REGEXMATCH": lambda s, p: bool(re.search(p, str(s))),
        "REGEXEXTRACT": lambda s, p: (re.search(p, str(s)).group(0) if re.search(p, str(s)) else ""),
        "REGEXREPLACE": lambda s, p, r: re.sub(p, r, str(s)),
        **_funcs_extended(),
    }


def _levenshtein(a: str, b: str) -> int:
    a, b = str(a), str(b)
    if len(a) < len(b):
        a, b = b, a
    prev = list(range(len(b) + 1))
    for i, ca in enumerate(a, 1):
        cur = [i]
        for j, cb in enumerate(b, 1):
            cur.append(min(prev[j] + 1, cur[j - 1] + 1, prev[j - 1] + (ca != cb)))
        prev = cur
    return prev[-1]


_ROMAN = [(1000, "M"), (900, "CM"), (500, "D"), (400, "CD"), (100, "C"), (90, "XC"),
          (50, "L"), (40, "XL"), (10, "X"), (9, "IX"), (5, "V"), (4, "IV"), (1, "I")]


def _roman(n: int) -> str:
    n, out = int(n), ""
    for v, sym in _ROMAN:
        while n >= v:
            out += sym
            n -= v
    return out


def _unroman(s: str) -> int:
    vals = {"I": 1, "V": 5, "X": 10, "L": 50, "C": 100, "D": 500, "M": 1000}
    s, total, prev = str(s).upper(), 0, 0
    for c in reversed(s):
        v = vals.get(c, 0)
        total += v if v >= prev else -v
        prev = max(prev, v)
    return total


def _dateparse(s: str) -> dt.date:
    for fmt in ("%Y-%m-%d", "%d/%m/%Y", "%m/%d/%Y", "%Y/%m/%d", "%d-%m-%Y"):
        try:
            return dt.datetime.strptime(str(s).strip(), fmt).date()
        except ValueError:
            continue
    raise ScriptError(f"Cannot parse date: {s}")


def _funcs_extended() -> Dict[str, Callable]:
    """Extended library: math, stats, strings, dates, lists, encoding,
    conversion, units, finance, geometry, colors, validation."""
    import statistics
    import uuid as _uuid
    import zlib
    import html as _html

    def _stat(fn, lst):
        vals = [_val(x) for x in lst]
        if not vals:
            return 0
        return fn(vals)

    def _pmt(rate, n, pv):
        rate, n, pv = float(rate), int(n), float(pv)
        if rate == 0:
            return pv / n
        return pv * rate / (1 - (1 + rate) ** -n)

    return {
        # ---- extra math
        "EXP": math.exp, "LN": math.log, "LOG10": math.log10, "LOG2": math.log2,
        "SIN": math.sin, "COS": math.cos, "TAN": math.tan,
        "ASIN": math.asin, "ACOS": math.acos, "ATAN": math.atan, "ATAN2": math.atan2,
        "SINH": math.sinh, "COSH": math.cosh, "TANH": math.tanh,
        "FLOOR": math.floor, "CEIL": math.ceil,
        "SIGN": lambda x: (x > 0) - (x < 0),
        "FRAC": lambda x: float(x) - int(float(x)),
        "GCD": lambda a, b: math.gcd(int(a), int(b)),
        "LCM": lambda a, b: abs(int(a) * int(b)) // math.gcd(int(a), int(b)) if a and b else 0,
        "FACT": lambda n: math.factorial(min(int(n), 170)),
        "COMB": lambda n, k: math.comb(min(int(n), 1000), int(k)),
        "PERM": lambda n, k: math.perm(min(int(n), 170), int(k)),
        "CLAMP": lambda x, lo, hi: max(lo, min(hi, x)),
        "LERP": lambda a, b, t: a + (b - a) * t,
        "MAPRANGE": lambda x, a1, b1, a2, b2: a2 + (x - a1) * (b2 - a2) / (b1 - a1),
        "DEG": math.degrees, "RAD": math.radians,
        "HYPOT": math.hypot, "CBRT": lambda x: float(x) ** (1 / 3),
        "PARSEINT": lambda s, base=10: int(str(s).strip(), int(base)),
        "PARSEFLOAT": lambda s: float(str(s).strip() or 0),
        # ---- statistics
        "MEAN": lambda lst: _stat(statistics.fmean, lst),
        "AVERAGE": lambda lst: _stat(statistics.fmean, lst),
        "MEDIAN": lambda lst: _stat(statistics.median, lst),
        "MODE": lambda lst: _stat(statistics.mode, lst),
        "STDEV": lambda lst: _stat(statistics.stdev, lst) if len(lst) > 1 else 0,
        "VARIANCE": lambda lst: _stat(statistics.variance, lst) if len(lst) > 1 else 0,
        "PERCENTILE": lambda lst, p: sorted(_val(x) for x in lst)[
            min(len(lst) - 1, max(0, int(round(float(p) / 100 * (len(lst) - 1)))))] if lst else 0,
        "ZSCORE": lambda x, lst: ((_val(x) - statistics.fmean([_val(v) for v in lst]))
                                  / (statistics.stdev([_val(v) for v in lst]) or 1)) if len(lst) > 1 else 0,
        # ---- random extras
        "CHOICE": lambda lst: random.choice(lst) if lst else "",
        "SHUFFLE": lambda lst: random.sample(list(lst), len(lst)),
        "SAMPLE": lambda lst, n: random.sample(list(lst), min(int(n), len(lst))),
        "UUID": lambda: str(_uuid.uuid4()),
        "RANDSTR": lambda n=8: "".join(random.choices(
            "ABCDEFGHJKLMNPQRSTUVWXYZabcdefghijkmnopqrstuvwxyz23456789", k=min(int(n), 256))),
        # ---- string extras
        "CAPITALIZE": lambda s: str(s).capitalize(),
        "TITLECASE": lambda s: str(s).title(),
        "SWAPCASE": lambda s: str(s).swapcase(),
        "PADLEFT": lambda s, n, c=" ": str(s).rjust(int(n), str(c)[:1] or " "),
        "PADRIGHT": lambda s, n, c=" ": str(s).ljust(int(n), str(c)[:1] or " "),
        "CENTER": lambda s, n, c=" ": str(s).center(int(n), str(c)[:1] or " "),
        "ZFILL": lambda s, n: str(s).zfill(int(n)),
        "LTRIM": lambda s: str(s).lstrip(), "RTRIM": lambda s: str(s).rstrip(),
        "SQUEEZE": lambda s: re.sub(r"\s+", " ", str(s)).strip(),
        "REVERSESTR": lambda s: str(s)[::-1],
        "COUNTSTR": lambda s, sub: str(s).count(str(sub)),
        "TRUNCATE": lambda s, n, suf="…": (str(s)[: int(n)] + str(suf)) if len(str(s)) > int(n) else str(s),
        "SLUG": lambda s: re.sub(r"[^a-z0-9]+", "-", str(s).lower()).strip("-"),
        "CAMELCASE": lambda s: (lambda w: (w[0].lower() + "".join(x.title() for x in w[1:])) if w else "")(
            re.split(r"[\s_\-]+", str(s).strip())),
        "SNAKECASE": lambda s: re.sub(r"[\s\-]+", "_", str(s).strip().lower()),
        "KEBABCASE": lambda s: re.sub(r"[\s_]+", "-", str(s).strip().lower()),
        "PASCALCASE": lambda s: "".join(w.title() for w in re.split(r"[\s_\-]+", str(s).strip())),
        "WRAP": lambda s, n=80: "\n".join(
            __import__("textwrap").wrap(str(s), int(n)) or [""]),
        "INDENT": lambda s, n=2: "\n".join(" " * int(n) + ln for ln in str(s).splitlines()),
        "QUOTE": lambda s: '"' + str(s).replace('"', '\\"') + '"',
        "UNQUOTE": lambda s: str(s)[1:-1] if len(str(s)) >= 2 and str(s)[0] == str(s)[-1] == '"' else str(s),
        "BETWEEN": lambda s, a, b: (str(s).split(str(a), 1)[1].split(str(b), 1)[0]
                                    if str(a) in str(s) and str(b) in str(s).split(str(a), 1)[1] else ""),
        "BEFORE": lambda s, sep: str(s).split(str(sep), 1)[0],
        "AFTER": lambda s, sep: str(s).split(str(sep), 1)[1] if str(sep) in str(s) else "",
        "REMOVE": lambda s, sub: str(s).replace(str(sub), ""),
        "INSERTSTR": lambda s, i, sub: str(s)[: int(i)] + str(sub) + str(s)[int(i):],
        "SPACES": lambda n: " " * int(n), "TAB": lambda: "\t", "NL": lambda: "\n",
        "LASTINSTR": lambda s, sub: str(s).rfind(str(sub)) + 1,
        "ISALPHA": lambda s: str(s).isalpha(), "ISDIGIT": lambda s: str(s).isdigit(),
        "ISALNUM": lambda s: str(s).isalnum(), "ISUPPER": lambda s: str(s).isupper(),
        "ISLOWER": lambda s: str(s).islower(), "ISSPACE": lambda s: str(s).isspace(),
        "LEVENSHTEIN": _levenshtein,
        "SIMILARITY": lambda a, b: round(1 - _levenshtein(a, b) / max(len(str(a)), len(str(b)), 1), 3),
        # ---- date extras
        "SECOND": lambda: dt.datetime.now().second,
        "DATEADD": lambda d, n: (_dateparse(d) + dt.timedelta(days=int(n))).isoformat(),
        "DATEDIFF": lambda a, b: (_dateparse(b) - _dateparse(a)).days,
        "DATEFORMAT": lambda d, fmt="%Y-%m-%d": _dateparse(d).strftime(str(fmt)),
        "DAYOFYEAR": lambda d=None: (_dateparse(d) if d else dt.date.today()).timetuple().tm_yday,
        "WEEKNUM": lambda d=None: (_dateparse(d) if d else dt.date.today()).isocalendar()[1],
        "ISLEAPYEAR": lambda y=None: __import__("calendar").isleap(int(y) if y else dt.date.today().year),
        "DAYSINMONTH": lambda y=None, m=None: __import__("calendar").monthrange(
            int(y) if y else dt.date.today().year, int(m) if m else dt.date.today().month)[1],
        "QUARTER": lambda d=None: ((_dateparse(d) if d else dt.date.today()).month - 1) // 3 + 1,
        "MONTHNAME": lambda m=None: __import__("calendar").month_name[
            int(m) if m else dt.date.today().month],
        "EPOCH2DATE": lambda t: dt.datetime.fromtimestamp(float(t)).strftime("%Y-%m-%d %H:%M:%S"),
        "DATE2EPOCH": lambda d: int(dt.datetime.combine(_dateparse(d), dt.time()).timestamp()),
        "AGE": lambda d: (dt.date.today() - _dateparse(d)).days // 365,
        # ---- list extras
        "SORTDESC": lambda lst: sorted(lst, reverse=True),
        "NTH": lambda lst, i: lst[int(i)] if -len(lst) <= int(i) < len(lst) else "",
        "INDEXOF": lambda lst, x: lst.index(x) if x in lst else -1,
        "FLATTEN": lambda lst: [y for x in lst for y in (x if isinstance(x, list) else [x])],
        "ZIP": lambda a, b: [list(p) for p in zip(a, b)],
        "FILL": lambda n, v: [v] * min(int(n), 100000),
        "SEQ": lambda a, b, s=1: list(range(int(a), int(b) + (1 if int(s) > 0 else -1), int(s))),
        "TAKE": lambda lst, n: list(lst)[: int(n)],
        "DROP": lambda lst, n: list(lst)[int(n):],
        "CHUNK": lambda lst, n: [list(lst)[i:i + int(n)] for i in range(0, len(lst), max(1, int(n)))],
        "INTERSECT": lambda a, b: [x for x in a if x in b],
        "UNIONL": lambda a, b: list(dict.fromkeys(list(a) + list(b))),
        "DIFFERENCE": lambda a, b: [x for x in a if x not in b],
        "PRODUCT": lambda lst: math.prod(_val(x) for x in lst),
        # ---- JSON extras
        "ISJSON": lambda s: (lambda: (json.loads(str(s)), True)[1])() if _try_json(s) else False,
        "KEYS": lambda o: list(o.keys()) if isinstance(o, dict) else [],
        "VALUES": lambda o: list(o.values()) if isinstance(o, dict) else [],
        "HASKEY": lambda o, k: str(k) in o if isinstance(o, dict) else False,
        # ---- encoding & crypto extras
        "MD5": lambda s: hashlib.md5(str(s).encode()).hexdigest(),
        "SHA1": lambda s: hashlib.sha1(str(s).encode()).hexdigest(),
        "SHA256": lambda s: hashlib.sha256(str(s).encode()).hexdigest(),
        "SHA512": lambda s: hashlib.sha512(str(s).encode()).hexdigest(),
        "CRC32": lambda s: format(zlib.crc32(str(s).encode()) & 0xFFFFFFFF, "08x"),
        "HEXENCODE": lambda s: str(s).encode().hex(),
        "HEXDECODE": lambda s: bytes.fromhex(str(s)).decode("utf-8", "replace"),
        "ROT13": lambda s: __import__("codecs").encode(str(s), "rot13"),
        "URLDECODE": lambda s: urllib.parse.unquote(str(s)),
        "HTMLESCAPE": lambda s: _html.escape(str(s)),
        "HTMLUNESCAPE": lambda s: _html.unescape(str(s)),
        # ---- regex extras
        "REGEXEXTRACTALL": lambda s, p: re.findall(p, str(s)),
        # ---- conversion & formatting
        "BIN": lambda n: bin(int(n))[2:], "OCT": lambda n: oct(int(n))[2:],
        "HEXN": lambda n: format(int(n), "x"),
        "TOBOOL": lambda x: str(x).strip().lower() in ("1", "true", "yes", "on", "y"),
        "CURRENCY": lambda x, sym="$": f"{sym}{float(x):,.2f}",
        "PERCENTF": lambda x, d=1: f"{float(x) * 100:.{int(d)}f}%",
        "COMMA": lambda x: f"{float(x):,.0f}" if float(x) == int(float(x)) else f"{float(x):,}",
        "FIXED": lambda x, d=2: f"{float(x):.{int(d)}f}",
        "ROMAN": _roman, "UNROMAN": _unroman,
        # ---- unit conversions
        "C2F": lambda c: float(c) * 9 / 5 + 32, "F2C": lambda f: (float(f) - 32) * 5 / 9,
        "KM2MI": lambda km: float(km) * 0.621371, "MI2KM": lambda mi: float(mi) * 1.609344,
        "KG2LB": lambda kg: float(kg) * 2.204623, "LB2KG": lambda lb: float(lb) * 0.453592,
        "M2FT": lambda m: float(m) * 3.28084, "FT2M": lambda ft: float(ft) * 0.3048,
        "L2GAL": lambda litre: float(litre) * 0.264172, "GAL2L": lambda g: float(g) * 3.785412,
        # ---- finance
        "PMT": _pmt,
        "FV": lambda pv, rate, n: float(pv) * (1 + float(rate)) ** int(n),
        "PV": lambda fv, rate, n: float(fv) / (1 + float(rate)) ** int(n),
        "COMPOUND": lambda p, rate, n, times=1: float(p) * (1 + float(rate) / int(times)) ** (int(times) * int(n)),
        # ---- geometry
        "AREACIRCLE": lambda r: math.pi * float(r) ** 2,
        "AREARECT": lambda w, h: float(w) * float(h),
        "AREATRI": lambda b, h: float(b) * float(h) / 2,
        "PYTHAG": lambda a, b: math.hypot(float(a), float(b)),
        # ---- colors
        "RGB2HEX": lambda r, g, b: "#{:02x}{:02x}{:02x}".format(
            max(0, min(255, int(r))), max(0, min(255, int(g))), max(0, min(255, int(b)))),
        "HEX2RGB": lambda h: [int(str(h).lstrip("#")[i:i + 2], 16) for i in (0, 2, 4)],
        # ---- validation
        "ISEMAIL": lambda s: bool(re.fullmatch(r"[^@\s]+@[^@\s]+\.[^@\s]+", str(s))),
        "ISURL": lambda s: bool(re.match(r"https?://[^\s]+\.[^\s]+", str(s))),
        "ISIP": lambda s: bool(re.fullmatch(
            r"(\d{1,3}\.){3}\d{1,3}", str(s))) and all(0 <= int(p) <= 255 for p in str(s).split(".")),
        "ISPHONE": lambda s: bool(re.fullmatch(r"\+?[\d\s\-().]{7,20}", str(s))),
        "ISUUID": lambda s: bool(re.fullmatch(
            r"[0-9a-fA-F]{8}-[0-9a-fA-F]{4}-[0-9a-fA-F]{4}-[0-9a-fA-F]{4}-[0-9a-fA-F]{12}", str(s))),
        "ISDATE": lambda s: _try_date(s),
        # ---- logic & misc
        "COALESCE": lambda *a: next((x for x in a if x not in ("", None, 0, False)), ""),
        "DEFAULTVAL": lambda x, d: d if x in ("", None) else x,
    }


def _try_json(s: Any) -> bool:
    try:
        json.loads(str(s))
        return True
    except Exception:
        return False


def _try_date(s: Any) -> bool:
    try:
        _dateparse(s)
        return True
    except Exception:
        return False


# --------------------------------------------------------- expression eval
_TOKEN_RE = re.compile(r'''("(?:[^"\\]|\\.)*")|(\bAND\b)|(\bOR\b)|(\bNOT\b)|(\bMOD\b)|(<>)|(&)|(\^)|(\\)''',
                       re.IGNORECASE)


def _to_python_expr(expr: str, cond: bool) -> str:
    """Translate BASIC operators to Python, protecting string literals."""
    def repl(m: re.Match) -> str:
        if m.group(1):
            return m.group(1)          # string literal untouched
        if m.group(2):
            return " and "
        if m.group(3):
            return " or "
        if m.group(4):
            return " not "
        if m.group(5):
            return " % "
        if m.group(6):
            return " != "
        if m.group(7):
            return " __CONCAT__ "      # placeholder, handled below
        if m.group(8):
            return " ** "
        if m.group(9):
            return " // "
        return m.group(0)

    out = _TOKEN_RE.sub(repl, expr)
    if cond:
        # single '=' → '==' (outside string literals; <>, <=, >=, == kept)
        parts = re.split(r'("(?:[^"\\]|\\.)*")', out)
        for i in range(0, len(parts), 2):
            parts[i] = re.sub(r"(?<![<>=!])=(?!=)", "==", parts[i])
        out = "".join(parts)
    return out


class _Env(dict):
    def __init__(self, vars: dict, funcs: dict):
        super().__init__()
        self._vars = vars
        self._funcs = funcs

    def __getitem__(self, key: str):
        ku = key.upper()
        if ku in self._funcs:
            return self._funcs[ku]
        if key in self._vars:
            return self._vars[key]
        if ku in self._vars:
            return self._vars[ku]
        raise ScriptError(f"Unknown variable or function: {key}")

    def __contains__(self, key):  # needed by eval for name lookups
        return True


def _concat(a, b):
    return str(a) + str(b)


class Interpreter:
    """Sandboxed SkillScript BASIC interpreter."""

    def __init__(self, code: str, context: Optional[dict] = None,
                 ai_runner: Optional[Callable[[str, str, str], str]] = None,
                 email_sender: Optional[Callable[[str, str, str], str]] = None):
        self.raw = code or ""
        self.context = dict(context or {})
        self.ai_runner = ai_runner
        self.email_sender = email_sender
        self.funcs = _funcs()
        self.vars: Dict[str, Any] = {}
        self.consts: set = set()
        self.output: List[str] = []
        self.trace: List[str] = []
        self.program_name = ""
        self.lines: List[tuple] = []       # (lineno, text)
        self.labels: Dict[str, int] = {}
        self.subs: Dict[str, int] = {}     # name -> index of SUB line
        self.on_error: Optional[str] = None  # label / "NEXT" / None
        self.steps = 0
        self.ai_calls = 0
        self.t0 = 0.0

    # ---------------------------------------------------------- parsing
    def _parse(self) -> None:
        idx = 0
        for rawline in self.raw.splitlines():
            s = rawline.strip()
            if not s or s.startswith("'"):
                continue
            s = re.sub(r"^\d+\s+", "", s)          # strip BASIC line numbers
            if s.upper().startswith("REM"):
                continue
            self.lines.append((idx, s))
            idx += 1
        # first pass: labels, subs, program name
        for i, (_, s) in enumerate(self.lines):
            u = s.upper()
            if u.startswith("PROGRAM "):
                self.program_name = s[8:].strip()
            elif u.startswith("SUB "):
                self.subs[s[4:].strip().upper()] = i
            elif re.fullmatch(r"[A-Za-z_][A-Za-z0-9_]*:", s):
                self.labels[s[:-1].upper()] = i

    # ------------------------------------------------------- expressions
    def eval(self, expr: str, cond: bool = False) -> Any:
        expr = expr.strip()
        py = _to_python_expr(expr, cond)
        py = py.replace("__CONCAT__", "+")
        env = _Env({**self._builtin_vars(), **self.vars},
                   {**self.funcs, "__CONCAT__": _concat})
        try:
            # string-concat friendliness: & already replaced by +; mixed
            # str+num handled by retry with str coercion
            return eval(compile(py, "<skillscript>", "eval"), {"__builtins__": {}}, env)
        except ScriptError:
            raise
        except TypeError:
            # retry treating + as string concatenation
            try:
                env2 = _Env({k: v for k, v in {**self._builtin_vars(), **self.vars}.items()},
                            self.funcs)
                parts = [self.eval(p) for p in self._split_top(expr, "&")]
                if len(parts) > 1:
                    return "".join(str(p) for p in parts)
                del env2
            except Exception:
                pass
            raise ScriptError(f"Type error in expression: {expr}")
        except ZeroDivisionError:
            raise ScriptError(f"Division by zero in: {expr}")
        except Exception as e:
            raise ScriptError(f"Bad expression: {expr} ({e})")

    @staticmethod
    def _split_top(s: str, sep: str) -> List[str]:
        parts, depth, cur, instr = [], 0, "", False
        i = 0
        while i < len(s):
            c = s[i]
            if c == '"':
                instr = not instr
            if not instr:
                if c in "([":
                    depth += 1
                elif c in ")]":
                    depth -= 1
                elif c == sep and depth == 0:
                    parts.append(cur)
                    cur = ""
                    i += 1
                    continue
            cur += c
            i += 1
        parts.append(cur)
        return parts

    def _builtin_vars(self) -> dict:
        now = dt.datetime.now()
        import socket
        return {
            "DATE$": now.strftime("%Y-%m-%d"), "TIME$": now.strftime("%H:%M:%S"),
            "NODE$": socket.gethostname(),
            "SKILL_NAME": self.context.get("skill_name", self.program_name),
            "TARGET": self.context.get("target", ""),
            "USER_INPUT": self.context.get("user_input", ""),
            "ERR": self.vars.get("ERR", 0), "ERRMSG": self.vars.get("ERRMSG", ""),
            "TRUE": True, "FALSE": False, "NULL": None, "PI": math.pi,
            **{k.upper(): v for k, v in self.context.items()
               if isinstance(v, (str, int, float, bool, list, dict))},
        }

    # ------------------------------------------------------------- output
    def _emit(self, text: str) -> None:
        self.output.append(str(text))
        if sum(len(x) for x in self.output) > MAX_OUTPUT:
            raise ScriptError("Output limit exceeded (512 KB)")

    def _log(self, msg: str) -> None:
        self.trace.append(msg[:500])
        if len(self.trace) > 500:
            del self.trace[:100]

    # ------------------------------------------------------------- limits
    def _tick(self) -> None:
        self.steps += 1
        if self.steps > MAX_STEPS:
            raise ScriptError(f"Step limit exceeded ({MAX_STEPS} statements)")
        if time.time() - self.t0 > MAX_SECONDS:
            raise ScriptError(f"Time limit exceeded ({MAX_SECONDS:.0f} s)")

    def _safe_path(self, p: str) -> Path:
        from .config import get_config
        cfg = get_config()
        roots = [DATA_DIR.resolve()]
        for key in ("files_dir", "images_dir"):
            try:
                roots.append(Path(cfg[key]).resolve())
            except Exception:
                pass
        target = Path(p)
        if not target.is_absolute():
            target = roots[0] / "skillfiles" / target
        target = target.resolve()
        for r in roots:
            try:
                target.relative_to(r)
                return target
            except ValueError:
                continue
        raise ScriptError(f"FILE access denied outside data folders: {p}")

    # ------------------------------------------------------- run the code
    def run(self) -> dict:
        self.t0 = time.time()
        self._parse()
        if not self.program_name:
            raise ScriptError("Missing PROGRAM statement")
        try:
            self._exec_block(0, len(self.lines))
        except _EndProgram:
            pass
        ms = int((time.time() - self.t0) * 1000)
        return {"ok": True, "name": self.program_name,
                "output": "".join(self.output), "trace": self.trace, "ms": ms,
                "steps": self.steps}

    # returns pointer control: None normal end
    def _exec_block(self, start: int, end: int) -> None:
        i = start
        gosub_stack: List[int] = []
        while i < end:
            self._tick()
            _, s = self.lines[i]
            try:
                nxt = self._exec_stmt(s, i, gosub_stack)
            except _EndProgram:
                raise
            except _FlowSignal:
                raise
            except ScriptError as e:
                i = self._handle_error(e, i, gosub_stack)
                if i is None:
                    raise
                continue
            i = nxt if nxt is not None else i + 1

    def _handle_error(self, e: ScriptError, i: int, gosub_stack) -> Optional[int]:
        self.vars["ERR"] = 1
        self.vars["ERRMSG"] = str(e)
        self._log(f"ERROR at stmt {i + 1}: {e}")
        if self.on_error == "NEXT":
            return i + 1
        if self.on_error and self.on_error in self.labels:
            lbl = self.on_error
            self.on_error = None          # avoid infinite error loops
            return self.labels[lbl]
        return None

    # --------------------------------------------------------- statements
    # First-word aliases resolved before dispatch (enterprise vocabulary)
    ALIASES = {
        "SAY": "PRINT", "ECHO": "PRINT", "WRITE": "PRINT", "WRITELN": "PRINT",
        "DISPLAY": "PRINT", "SHOW": "PRINT", "TELL": "PRINT", "EMIT": "PRINT",
        "WAIT": "SLEEP", "PAUSE": "SLEEP", "DELAY": "SLEEP",
        "SET": "LET", "VAR": "LET", "DECLARE": "LET", "GLOBAL": "LET", "LOCAL": "LET",
        "DEFINE": "CONST",
        "HALT": "END", "QUIT": "END", "FINISH": "END", "TERMINATE": "END",
        "FAIL": "THROW", "PANIC": "THROW", "RAISE": "THROW",
        "REQUIRE": "ASSERT", "VERIFY": "ASSERT", "EXPECT": "ASSERT", "CHECK": "ASSERT",
        "DEBUG": "LOG", "TRACE": "LOG", "NOTE": "LOG", "INFO": "LOG", "WARN": "LOG",
        "JUMP": "GOTO",
        "PERFORM": "CALL", "INVOKE": "CALL",
        "FETCH": "HTTP",
        "QUERY": "ASK", "LLM": "ASK", "GPT": "ASK", "AIASK": "ASK",
    }

    def _exec_stmt(self, s: str, i: int, gosub_stack: List[int]) -> Optional[int]:
        # resolve first-word aliases
        first = s.split(None, 1)[0].upper() if s.split() else ""
        if first in self.ALIASES:
            rest = s.split(None, 1)[1] if len(s.split(None, 1)) > 1 else ""
            s = (self.ALIASES[first] + (" " + rest if rest else "")).strip()
        u = s.upper()

        # --- structure / no-ops
        if u.startswith("PROGRAM ") or u == "END PROGRAM" or re.fullmatch(r"[A-Za-z_][A-Za-z0-9_]*:", s):
            return None
        if u in ("NOP", "BEEP"):
            return None
        if u in ("END", "STOP"):
            raise _EndProgram()
        if u.startswith("ABORT"):
            msg = self.eval(s[5:]) if s[5:].strip() else "aborted"
            raise ScriptError(f"ABORT: {msg}")
        if u.startswith("THROW "):
            raise ScriptError(str(self.eval(s[6:])))

        # --- output helpers
        if u == "CLS":
            self.output.clear()
            return None
        if u == "HR" or u.startswith("HR "):
            ch = str(self.eval(s[2:])) if s[2:].strip() else "-"
            self._emit((ch[:1] or "-") * 40 + "\n")
            return None
        if u.startswith("TITLE "):
            self._emit("# " + str(self.eval(s[6:])) + "\n")
            return None
        if u.startswith("HEADER "):
            self._emit("## " + str(self.eval(s[7:])) + "\n")
            return None
        if u.startswith("BULLET "):
            self._emit("* " + str(self.eval(s[7:])) + "\n")
            return None

        # --- variable helpers
        if u.startswith("SWAP "):
            parts = [p.strip() for p in self._split_top(s[5:], ",")]
            if len(parts) != 2:
                raise ScriptError("SWAP needs: SWAP a, b")
            a, b = parts
            va, vb = self.vars.get(a), self.vars.get(b)
            if a not in self.vars or b not in self.vars:
                raise ScriptError(f"SWAP: unknown variable {a if a not in self.vars else b}")
            self.vars[a], self.vars[b] = vb, va
            return None
        if u.startswith("INC ") or u.startswith("DEC "):
            sign = 1 if u.startswith("INC") else -1
            parts = [p.strip() for p in self._split_top(s[4:], ",")]
            var = parts[0]
            amount = self.eval(parts[1]) if len(parts) > 1 else 1
            self.vars[var] = _val(self.vars.get(var, 0)) + sign * amount
            return None
        if u.startswith("PUSH "):
            parts = self._split_top(s[5:], ",")
            if len(parts) < 2:
                raise ScriptError("PUSH needs: PUSH list, value")
            name = parts[0].strip()
            lst = self.vars.get(name)
            if not isinstance(lst, list):
                raise ScriptError(f"PUSH: {name} is not a list")
            lst.append(self.eval(",".join(parts[1:])))
            return None
        if u.startswith("POP "):
            m = re.match(r"POP\s+([A-Za-z_][A-Za-z0-9_]*)\s*=\s*([A-Za-z_][A-Za-z0-9_]*)\s*$", s, re.IGNORECASE)
            if not m:
                raise ScriptError("POP needs: POP var = list")
            lst = self.vars.get(m.group(2))
            if not isinstance(lst, list):
                raise ScriptError(f"POP: {m.group(2)} is not a list")
            self.vars[m.group(1)] = lst.pop() if lst else ""
            return None
        if u.startswith("RANDOMIZE"):
            seed = s[9:].strip()
            random.seed(self.eval(seed) if seed else None)
            return None

        # --- SUB definitions are skipped in normal flow
        if u.startswith("SUB "):
            return self._find_forward(i, "END SUB") + 1
        if u == "END SUB":
            raise _ReturnSignal()
        if u.startswith("CALL "):
            name = s[5:].strip().upper()
            if name not in self.subs:
                raise ScriptError(f"Unknown SUB: {name}")
            try:
                self._exec_block(self.subs[name] + 1, len(self.lines))
            except _ReturnSignal:
                pass
            return None

        # --- flow
        if u.startswith("GOTO "):
            lbl = s[5:].strip().upper()
            if lbl not in self.labels:
                raise ScriptError(f"Unknown label: {lbl}")
            return self.labels[lbl]
        if u.startswith("GOSUB "):
            lbl = s[6:].strip().upper()
            if lbl not in self.labels:
                raise ScriptError(f"Unknown label: {lbl}")
            gosub_stack.append(i + 1)
            return self.labels[lbl]
        if u == "RETURN":
            if not gosub_stack:
                raise ScriptError("RETURN without GOSUB")
            return gosub_stack.pop()
        if u in ("BREAK", "EXIT FOR", "EXIT WHILE", "EXIT DO", "EXIT LOOP"):
            raise _BreakSignal()
        if u == "CONTINUE":
            raise _ContinueSignal()

        # --- error handling policy
        if u.startswith("ON ERROR"):
            rest = s[8:].strip().upper()
            if rest.startswith("GOTO"):
                self.on_error = rest[4:].strip()
            elif rest.startswith("RESUME NEXT"):
                self.on_error = "NEXT"
            else:
                self.on_error = None
            return None

        # --- IF / block structures
        if u.startswith("IF "):
            return self._exec_if(s, i, gosub_stack)
        if u.startswith(("ELSEIF ", "ELSE IF ")) or u == "ELSE":
            return self._find_forward(i, "ENDIF", "END IF") + 1  # skip taken-branch leftovers
        if u in ("ENDIF", "END IF"):
            return None
        if u.startswith("SELECT CASE"):
            return self._exec_select(s, i, gosub_stack)
        if u.startswith("CASE"):
            return self._find_forward(i, "END SELECT") + 1
        if u == "END SELECT":
            return None
        if u.startswith("FOR "):
            return self._exec_for(s, i, gosub_stack)
        if u.startswith("NEXT"):
            return None
        if u.startswith("WHILE "):
            return self._exec_while(s, i, gosub_stack)
        if u == "WEND":
            return None
        if u == "DO":
            return self._exec_do(i, gosub_stack)
        if u.startswith("LOOP"):
            return None

        # --- IO & effects
        if u.startswith("PRINT") or u.startswith("OUTPUT"):
            rest = s[5:] if u.startswith("PRINT") else s[6:]
            rest = rest.strip()
            if not rest:
                self._emit("\n")
                return None
            parts = self._split_top(rest, ";")
            text = "".join(str(self.eval(p)) for p in parts if p.strip())
            self._emit(text + ("" if rest.endswith(";") else "\n"))
            return None
        if u.startswith("LOG "):
            self._log(str(self.eval(s[4:])))
            return None
        if u.startswith("INPUT "):
            rest = s[6:].strip()
            var = rest.split(",")[0].strip()
            self.vars[var] = self.context.get(var, self.context.get(var.lower(), ""))
            return None
        if u.startswith("SLEEP"):
            sec = min(float(self.eval(s[5:] or "1")), MAX_SLEEP)
            time.sleep(max(0.0, sec))
            return None
        if u.startswith("ASSERT "):
            parts = self._split_top(s[7:], ",")
            if not self.eval(parts[0], cond=True):
                msg = str(self.eval(parts[1])) if len(parts) > 1 else f"assertion failed: {parts[0].strip()}"
                raise ScriptError(f"ASSERT: {msg}")
            return None

        # --- AI
        if u.startswith("ASK "):
            return self._exec_ask(s)
        if u.startswith(("SUMMARIZE ", "TRANSLATE ", "SENTIMENT ", "CLASSIFY ",
                         "EXTRACT ", "REWRITE ", "EXPAND ", "OUTLINE ", "PROOFREAD ",
                         "KEYWORDS ", "TITLEGEN ", "REPLYGEN ", "CODEGEN ", "SQLGEN ")):
            return self._exec_ai_helper(s)
        if u.startswith("IMAGE "):
            return self._exec_image(s)

        # --- document / chart generation
        if u.startswith("GENERATE "):
            return self._exec_generate(s)
        if u.startswith("CHART "):
            return self._exec_chart(s)

        # --- HTTP
        if u.startswith("HTTP "):
            return self._exec_http(s)

        # --- FILE
        if u.startswith("FILE "):
            return self._exec_file(s)

        # --- MEM
        if u.startswith("MEM "):
            return self._exec_mem(s)

        # --- EMAIL
        if u.startswith("EMAIL "):
            return self._exec_email(s)

        # --- declarations & assignment
        if u.startswith("CONST "):
            name, expr = self._split_assign(s[6:])
            self.vars[name] = self.eval(expr)
            self.consts.add(name.upper())
            return None
        if u.startswith("DIM "):
            m = re.fullmatch(r"DIM\s+([A-Za-z_][A-Za-z0-9_]*)\s*\(\s*(.+)\s*\)", s, re.IGNORECASE)
            if not m:
                raise ScriptError(f"Bad DIM: {s}")
            self.vars[m.group(1)] = [""] * int(self.eval(m.group(2)))
            return None
        if u.startswith("LET "):
            s = s[4:]
        if "=" in s:
            name, expr = self._split_assign(s)
            if name.upper() in self.consts:
                raise ScriptError(f"Cannot reassign CONST {name}")
            # array element assignment a(i) = x  or a[i] = x
            m = re.fullmatch(r"([A-Za-z_][A-Za-z0-9_]*)\s*[\(\[]\s*(.+)\s*[\)\]]", name)
            if m:
                arr = self.vars.get(m.group(1))
                if not isinstance(arr, list):
                    raise ScriptError(f"{m.group(1)} is not an array")
                arr[int(self.eval(m.group(2)))] = self.eval(expr)
            else:
                if not re.fullmatch(r"[A-Za-z_][A-Za-z0-9_$]*", name):
                    raise ScriptError(f"Bad variable name: {name}")
                self.vars[name] = self.eval(expr)
            return None

        raise ScriptError(f"Unknown statement: {s}")

    @staticmethod
    def _split_assign(s: str) -> tuple:
        # split on first top-level '='
        depth = 0
        instr = False
        for idx, c in enumerate(s):
            if c == '"':
                instr = not instr
            if instr:
                continue
            if c in "([":
                depth += 1
            elif c in ")]":
                depth -= 1
            elif c == "=" and depth == 0:
                return s[:idx].strip(), s[idx + 1:].strip()
        raise ScriptError(f"Expected '=' in: {s}")

    # ----------------------------------------------------- block helpers
    def _find_forward(self, i: int, *terminators: str) -> int:
        """Find matching terminator, respecting nesting of the same opener."""
        opener = self.lines[i][1].split()[0].upper()
        openers = {"IF": ("IF ",), "FOR": ("FOR ",), "WHILE": ("WHILE ",),
                   "DO": ("DO",), "SELECT": ("SELECT CASE",), "SUB": ("SUB ",)}
        opens = openers.get(opener, (opener + " ",))
        depth = 0
        terms = tuple(t.upper() for t in terminators)
        for j in range(i + 1, len(self.lines)):
            u = self.lines[j][1].upper()
            if any(u.startswith(o) for o in opens) and not (opener == "IF" and "THEN" in u and not u.endswith("THEN")):
                depth += 1
            elif any(u == t or u.startswith(t + " ") for t in terms):
                if depth == 0:
                    return j
                depth -= 1
        raise ScriptError(f"Missing {'/'.join(terminators)} for {opener} (statement {i + 1})")

    def _exec_if(self, s: str, i: int, gosub_stack) -> Optional[int]:
        m = re.match(r"IF\s+(.*?)\s+THEN\s*(.*)$", s, re.IGNORECASE | re.DOTALL)
        if not m:
            raise ScriptError(f"Bad IF: {s}")
        cond, inline = m.group(1), m.group(2).strip()
        if inline:  # inline form: IF x THEN stmt
            if self.eval(cond, cond=True):
                return self._exec_stmt(inline, i, gosub_stack)
            return None
        endif = self._find_forward(i, "ENDIF", "END IF")
        # collect branch boundaries
        branches = [(i, cond)]
        depth = 0
        for j in range(i + 1, endif):
            u = self.lines[j][1].upper()
            if u.startswith("IF ") and u.endswith("THEN"):
                depth += 1
            elif u in ("ENDIF", "END IF"):
                depth -= 1
            elif depth == 0 and (u.startswith(("ELSEIF ", "ELSE IF "))):
                c = re.match(r"ELSE\s?IF\s+(.*?)\s+THEN\s*$", self.lines[j][1], re.IGNORECASE).group(1)
                branches.append((j, c))
            elif depth == 0 and u == "ELSE":
                branches.append((j, None))
        branches.append((endif, "__END__"))
        for b in range(len(branches) - 1):
            start, c = branches[b]
            if c == "__END__":
                break
            if c is None or self.eval(c, cond=True):
                self._exec_block(start + 1, branches[b + 1][0])
                break
        return endif + 1

    def _exec_select(self, s: str, i: int, gosub_stack) -> int:
        value = self.eval(s[len("SELECT CASE"):])
        endsel = self._find_forward(i, "END SELECT")
        cases = []
        for j in range(i + 1, endsel):
            u = self.lines[j][1].upper()
            if u.startswith("CASE"):
                cases.append(j)
        cases.append(endsel)
        for c in range(len(cases) - 1):
            line = self.lines[cases[c]][1]
            rest = line[4:].strip()
            if rest.upper() == "ELSE":
                match = True
            else:
                match = any(self.eval(p) == value for p in self._split_top(rest, ","))
            if match:
                self._exec_block(cases[c] + 1, cases[c + 1])
                break
        return endsel + 1

    def _exec_for(self, s: str, i: int, gosub_stack) -> int:
        m = re.match(r"FOR\s+([A-Za-z_][A-Za-z0-9_]*)\s*=\s*(.+?)\s+TO\s+(.+?)(?:\s+STEP\s+(.+))?$",
                     s, re.IGNORECASE)
        if not m:
            raise ScriptError(f"Bad FOR: {s}")
        var, a, b = m.group(1), self.eval(m.group(2)), self.eval(m.group(3))
        step = self.eval(m.group(4)) if m.group(4) else 1
        if step == 0:
            raise ScriptError("FOR STEP cannot be 0")
        nexti = self._find_forward(i, "NEXT")
        cur = a
        while (step > 0 and cur <= b) or (step < 0 and cur >= b):
            self._tick()
            self.vars[var] = cur
            try:
                self._exec_block(i + 1, nexti)
            except _BreakSignal:
                break
            except _ContinueSignal:
                pass
            cur += step
        return nexti + 1

    def _exec_while(self, s: str, i: int, gosub_stack) -> int:
        cond = s[6:]
        wend = self._find_forward(i, "WEND")
        while self.eval(cond, cond=True):
            self._tick()
            try:
                self._exec_block(i + 1, wend)
            except _BreakSignal:
                break
            except _ContinueSignal:
                pass
        return wend + 1

    def _exec_do(self, i: int, gosub_stack) -> int:
        loop = self._find_forward(i, "LOOP")
        loopline = self.lines[loop][1]
        m = re.match(r"LOOP\s+(WHILE|UNTIL)\s+(.*)$", loopline, re.IGNORECASE)
        while True:
            self._tick()
            try:
                self._exec_block(i + 1, loop)
            except _BreakSignal:
                break
            except _ContinueSignal:
                pass
            if not m:
                raise ScriptError("DO without LOOP WHILE/UNTIL condition")
            ok = self.eval(m.group(2), cond=True)
            if (m.group(1).upper() == "WHILE" and not ok) or (m.group(1).upper() == "UNTIL" and ok):
                break
        return loop + 1

    # -------------------------------------------------------- effect ops
    def _exec_ask(self, s: str) -> None:
        m = re.match(r"ASK\s+([A-Za-z_][A-Za-z0-9_]*)\s*=\s*(.*)$", s, re.IGNORECASE | re.DOTALL)
        if not m:
            raise ScriptError(f"Bad ASK (use: ASK var = prompt [SYSTEM expr] [USING \"codex\"]): {s}")
        var, rest = m.group(1), m.group(2)
        agent = "codex"
        um = re.search(r'\s+USING\s+("(?:[^"]*)")\s*$', rest, re.IGNORECASE)
        if um:
            agent = um.group(1).strip('"').lower()
            rest = rest[: um.start()]
        system = ""
        sm = re.search(r"\s+SYSTEM\s+(.+)$", rest, re.IGNORECASE)
        if sm:
            system = str(self.eval(sm.group(1)))
            rest = rest[: sm.start()]
        prompt = str(self.eval(rest))
        self.ai_calls += 1
        if self.ai_calls > MAX_AI_CALLS:
            raise ScriptError(f"AI call limit exceeded ({MAX_AI_CALLS} per run)")
        if not self.ai_runner:
            raise ScriptError("ASK unavailable in this context (no AI runner)")
        self._log(f"ASK {agent}: {prompt[:120]}")
        self.vars[var] = self.ai_runner(agent, prompt, system)

    def _ai(self, prompt: str, system: str = "", agent: str = "codex") -> str:
        self.ai_calls += 1
        if self.ai_calls > MAX_AI_CALLS:
            raise ScriptError(f"AI call limit exceeded ({MAX_AI_CALLS} per run)")
        if not self.ai_runner:
            raise ScriptError("AI statements unavailable in this context (no AI runner)")
        return self.ai_runner(agent, prompt, system)

    # AI convenience statements — sugar over ASK with task-tuned prompts
    _AI_HELPERS = {
        "SUMMARIZE": ("Summarize the following text concisely{opt}. Reply with the summary only.",
                      "WORDS", " in about {v} words"),
        "TRANSLATE": ("Translate the following text{opt}. Reply with the translation only.",
                      "TO", " to {v}"),
        "SENTIMENT": ("Analyze the sentiment of the following text. Reply with exactly one word: "
                      "positive, negative or neutral.", None, ""),
        "CLASSIFY": ("Classify the following text into exactly one of these categories{opt}. "
                     "Reply with the category name only.", "INTO", ": {v}"),
        "EXTRACT": ("Extract the following fields from the text{opt}. Reply as compact JSON only.",
                    "FIELDS", ": {v}"),
        "REWRITE": ("Rewrite the following text{opt}. Keep the meaning; reply with the rewritten "
                    "text only.", "STYLE", " in this style: {v}"),
        "EXPAND": ("Expand the following outline/notes into full prose{opt}. Reply with the text "
                   "only.", "WORDS", " of about {v} words"),
        "OUTLINE": ("Create a structured outline (markdown bullets) for the following topic/text"
                    "{opt}. Reply with the outline only.", "DEPTH", " with {v} levels"),
        "PROOFREAD": ("Proofread and correct grammar/spelling of the following text. Reply with "
                      "the corrected text only.", None, ""),
        "KEYWORDS": ("Extract the most important keywords from the following text{opt}. Reply as "
                     "a comma-separated list only.", "TOP", " (top {v})"),
        "TITLEGEN": ("Generate a compelling title for the following text{opt}. Reply with the "
                     "title only.", "STYLE", " in this style: {v}"),
        "REPLYGEN": ("Write a professional reply to the following message{opt}. Reply with the "
                     "response text only.", "TONE", " with a {v} tone"),
        "CODEGEN": ("Write code for the following task{opt}. Reply with code only, no explanations.",
                    "LANG", " in {v}"),
        "SQLGEN": ("Write a SQL query for the following request{opt}. Reply with SQL only.",
                   "DIALECT", " for {v}"),
    }

    def _exec_ai_helper(self, s: str) -> None:
        kw = s.split(None, 1)[0].upper()
        tmpl, modkw, modfmt = self._AI_HELPERS[kw]
        m = re.match(rf"{kw}\s+([A-Za-z_][A-Za-z0-9_]*)\s*=\s*(.*)$", s, re.IGNORECASE | re.DOTALL)
        if not m:
            raise ScriptError(f"Bad {kw} (use: {kw} var = text_expr"
                              + (f" [{modkw} expr]" if modkw else "") + ")")
        var, rest = m.group(1), m.group(2)
        opt = ""
        if modkw:
            mm = re.search(rf"\s+{modkw}\s+(.+)$", rest, re.IGNORECASE | re.DOTALL)
            if mm:
                opt = modfmt.replace("{v}", str(self.eval(mm.group(1))))
                rest = rest[: mm.start()]
        text = str(self.eval(rest))
        prompt = tmpl.replace("{opt}", opt) + "\n\n---\n" + text
        self._log(f"{kw}: {text[:100]}")
        self.vars[var] = self._ai(prompt).strip()

    def _exec_image(self, s: str) -> None:
        # IMAGE path_expr, prompt_expr — AI image generation into the images dir
        parts = self._split_top(s[6:], ",")
        if len(parts) < 2:
            raise ScriptError('Bad IMAGE (use: IMAGE "name.png", "description of the image")')
        path = self._safe_path(str(self.eval(parts[0])))
        prompt = str(self.eval(",".join(parts[1:])))
        self._log(f"IMAGE {path.name}: {prompt[:100]}")
        result = self._ai(
            f"Generate an image and save it EXACTLY to this file path: {path}\n"
            f"Image description: {prompt}\n"
            "Use any available image-generation capability or draw it as SVG/PNG with a script. "
            "Reply only with: OK <path>  when the file exists.")
        self.vars["IMAGE_RESULT"] = result.strip()

    # ---------------------------------------------------- file generation
    def _exec_generate(self, s: str) -> None:
        m = re.match(r"GENERATE\s+(TXT|TEXT|MD|MARKDOWN|HTML|CSV|JSON|XML|PDF|DOCX|XLSX|PPTX)\s+(.*)$",
                     s, re.IGNORECASE | re.DOTALL)
        if not m:
            raise ScriptError('Bad GENERATE (use: GENERATE PDF|DOCX|XLSX|PPTX|HTML|TXT|MD|CSV|JSON|XML '
                              '"file.ext", content)')
        fmt = m.group(1).upper()
        parts = self._split_top(m.group(2), ",")
        if len(parts) < 2:
            raise ScriptError(f"GENERATE {fmt} needs: path, content")
        path = self._safe_path(str(self.eval(parts[0])))
        content = self.eval(",".join(parts[1:]))
        path.parent.mkdir(parents=True, exist_ok=True)
        self._log(f"GENERATE {fmt} {path.name}")
        if fmt in ("TXT", "TEXT", "MD", "MARKDOWN", "HTML", "XML"):
            path.write_text(str(content), encoding="utf-8")
        elif fmt == "JSON":
            path.write_text(content if isinstance(content, str) and _try_json(content)
                            else json.dumps(content, ensure_ascii=False, indent=1), encoding="utf-8")
        elif fmt == "CSV":
            import csv as _csv
            with open(path, "w", newline="", encoding="utf-8-sig") as f:
                w = _csv.writer(f)
                if isinstance(content, list) and content and isinstance(content[0], list):
                    w.writerows(content)
                elif isinstance(content, list):
                    for item in content:
                        w.writerow([item])
                else:
                    f.write(str(content))
        elif fmt == "PDF":
            self._gen_pdf(path, str(content))
        elif fmt == "DOCX":
            self._gen_docx(path, str(content))
        elif fmt == "XLSX":
            self._gen_xlsx(path, content)
        elif fmt == "PPTX":
            self._gen_pptx(path, str(content))
        self.vars["LAST_FILE"] = str(path)

    @staticmethod
    def _gen_pdf(path: Path, text: str) -> None:
        """Minimal dependency-free single/multi-page text PDF."""
        def pdf_escape(t: str) -> str:
            return t.replace("\\", r"\\").replace("(", r"\(").replace(")", r"\)")
        lines = []
        for para in text.splitlines():
            while len(para) > 95:
                lines.append(para[:95])
                para = para[95:]
            lines.append(para)
        pages = [lines[i:i + 54] for i in range(0, max(len(lines), 1), 54)]
        objs: List[bytes] = []
        page_ids = []
        content_ids = []
        nobj = 3  # 1=catalog 2=pages, then per page: page obj + content obj
        for pg in pages:
            page_ids.append(nobj)
            content_ids.append(nobj + 1)
            nobj += 2
        font_id = nobj
        header = b"%PDF-1.4\n"
        body = b""
        offsets = {}
        def add(oid: int, data: bytes):
            nonlocal body
            offsets[oid] = len(header) + len(body)
            body += f"{oid} 0 obj\n".encode() + data + b"\nendobj\n"
        kids = " ".join(f"{pid} 0 R" for pid in page_ids)
        add(1, b"<< /Type /Catalog /Pages 2 0 R >>")
        add(2, f"<< /Type /Pages /Kids [{kids}] /Count {len(pages)} >>".encode())
        for pid, cid, pg in zip(page_ids, content_ids, pages):
            add(pid, (f"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
                      f"/Contents {cid} 0 R /Resources << /Font << /F1 {font_id} 0 R >> >> >>").encode())
            stream = "BT /F1 11 Tf 40 750 Td 13 TL\n"
            for ln in pg:
                stream += f"({pdf_escape(ln)}) Tj T*\n"
            stream += "ET"
            sb = stream.encode("latin-1", "replace")
            add(cid, f"<< /Length {len(sb)} >>\nstream\n".encode() + sb + b"\nendstream")
        add(font_id, b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>")
        xref_pos = len(header) + len(body)
        n = font_id + 1
        xref = f"xref\n0 {n}\n0000000000 65535 f \n".encode()
        for oid in range(1, n):
            xref += f"{offsets[oid]:010d} 00000 n \n".encode()
        trailer = (f"trailer\n<< /Size {n} /Root 1 0 R >>\nstartxref\n{xref_pos}\n%%EOF").encode()
        path.write_bytes(header + body + xref + trailer)

    @staticmethod
    def _gen_docx(path: Path, text: str) -> None:
        try:
            import docx  # python-docx if installed
            d = docx.Document()
            for para in text.splitlines():
                if para.startswith("# "):
                    d.add_heading(para[2:], level=1)
                elif para.startswith("## "):
                    d.add_heading(para[3:], level=2)
                else:
                    d.add_paragraph(para)
            d.save(str(path))
            return
        except ImportError:
            pass
        # dependency-free minimal DOCX (zip + document.xml)
        import zipfile
        from xml.sax.saxutils import escape as x
        paras = "".join(
            f'<w:p><w:pPr><w:pStyle w:val="{"Heading1" if p.startswith("# ") else "Heading2" if p.startswith("## ") else "Normal"}"/></w:pPr>'
            f'<w:r><w:t xml:space="preserve">{x(p.lstrip("# "))}</w:t></w:r></w:p>'
            for p in text.splitlines() or [""])
        doc = ('<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
               '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
               f'<w:body>{paras}</w:body></w:document>')
        with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as z:
            z.writestr("[Content_Types].xml",
                       '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
                       '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
                       '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
                       '<Default Extension="xml" ContentType="application/xml"/>'
                       '<Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/></Types>')
            z.writestr("_rels/.rels",
                       '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
                       '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
                       '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/></Relationships>')
            z.writestr("word/document.xml", doc)

    @staticmethod
    def _gen_xlsx(path: Path, content: Any) -> None:
        rows = content if isinstance(content, list) else [[content]]
        rows = [r if isinstance(r, list) else [r] for r in rows]
        try:
            import openpyxl
            wb = openpyxl.Workbook()
            ws = wb.active
            for r in rows:
                ws.append(r)
            wb.save(str(path))
            return
        except ImportError:
            pass
        # dependency-free minimal XLSX with inline strings
        import zipfile
        from xml.sax.saxutils import escape as x

        def cell(v, cref):
            if isinstance(v, (int, float)) and not isinstance(v, bool):
                return f'<c r="{cref}"><v>{v}</v></c>'
            return f'<c r="{cref}" t="inlineStr"><is><t xml:space="preserve">{x(str(v))}</t></is></c>'

        def colletter(n):
            out = ""
            while n >= 0:
                out = chr(65 + n % 26) + out
                n = n // 26 - 1
            return out
        rows_xml = "".join(
            f'<row r="{ri + 1}">' + "".join(cell(v, f"{colletter(ci)}{ri + 1}") for ci, v in enumerate(row)) + "</row>"
            for ri, row in enumerate(rows))
        with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as z:
            z.writestr("[Content_Types].xml",
                       '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
                       '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
                       '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
                       '<Default Extension="xml" ContentType="application/xml"/>'
                       '<Override PartName="/xl/workbook.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet.main+xml"/>'
                       '<Override PartName="/xl/worksheets/sheet1.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.worksheet+xml"/></Types>')
            z.writestr("_rels/.rels",
                       '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
                       '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
                       '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="xl/workbook.xml"/></Relationships>')
            z.writestr("xl/workbook.xml",
                       '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
                       '<workbook xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" '
                       'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
                       '<sheets><sheet name="Sheet1" sheetId="1" r:id="rId1"/></sheets></workbook>')
            z.writestr("xl/_rels/workbook.xml.rels",
                       '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
                       '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
                       '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/worksheet" Target="worksheets/sheet1.xml"/></Relationships>')
            z.writestr("xl/worksheets/sheet1.xml",
                       '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
                       '<worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">'
                       f'<sheetData>{rows_xml}</sheetData></worksheet>')

    def _gen_pptx(self, path: Path, text: str) -> None:
        try:
            from pptx import Presentation  # python-pptx if installed
            prs = Presentation()
            for block in text.split("\n---\n"):
                lines = block.strip().splitlines() or [""]
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = lines[0].lstrip("# ")
                if len(lines) > 1:
                    slide.placeholders[1].text = "\n".join(ln.lstrip("*- ") for ln in lines[1:])
            prs.save(str(path))
            return
        except ImportError:
            raise ScriptError("GENERATE PPTX requires the python-pptx package on the server "
                              "(pip install python-pptx) — or GENERATE PDF/DOCX/HTML instead")

    # ------------------------------------------------------------ charts
    def _exec_chart(self, s: str) -> None:
        # CHART "file.svg", "bar"|"line"|"pie", labels, values [, title]
        parts = self._split_top(s[6:], ",")
        if len(parts) < 4:
            raise ScriptError('Bad CHART (use: CHART "chart.svg", "bar"|"line"|"pie", labels, values [, title])')
        path = self._safe_path(str(self.eval(parts[0])))
        kind = str(self.eval(parts[1])).lower()
        labels = [str(x) for x in self.eval(parts[2])]
        values = [_val(x) for x in self.eval(parts[3])]
        title = str(self.eval(parts[4])) if len(parts) > 4 else ""
        if kind not in ("bar", "line", "pie"):
            raise ScriptError('CHART type must be "bar", "line" or "pie"')
        if len(labels) != len(values) or not values:
            raise ScriptError("CHART: labels and values must be equal-length non-empty lists")
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_text(self._svg_chart(kind, labels, values, title), encoding="utf-8")
        self._log(f"CHART {kind} {path.name} ({len(values)} points)")
        self.vars["LAST_FILE"] = str(path)

    @staticmethod
    def _svg_chart(kind: str, labels: List[str], values: List[float], title: str) -> str:
        from xml.sax.saxutils import escape as x
        W, H, PAD = 640, 400, 50
        colors = ["#3b82f6", "#22c55e", "#eab308", "#ef4444", "#a78bfa", "#22d3ee",
                  "#fb923c", "#f472b6", "#84cc16", "#64748b"]
        out = [f'<svg xmlns="http://www.w3.org/2000/svg" width="{W}" height="{H}" '
               f'viewBox="0 0 {W} {H}" font-family="Segoe UI,sans-serif">',
               f'<rect width="{W}" height="{H}" fill="#ffffff"/>']
        if title:
            out.append(f'<text x="{W/2}" y="26" text-anchor="middle" font-size="17" '
                       f'font-weight="bold" fill="#111">{x(title)}</text>')
        vmax = max(max(values), 1e-9)
        if kind == "pie":
            cx, cy, r = W / 2, H / 2 + 10, min(W, H) / 2 - PAD - 20
            total = sum(values) or 1
            ang = -90.0
            for i, (lb, v) in enumerate(zip(labels, values)):
                sweep = 360 * v / total
                a1, a2 = math.radians(ang), math.radians(ang + sweep)
                x1, y1 = cx + r * math.cos(a1), cy + r * math.sin(a1)
                x2, y2 = cx + r * math.cos(a2), cy + r * math.sin(a2)
                large = 1 if sweep > 180 else 0
                out.append(f'<path d="M{cx},{cy} L{x1:.1f},{y1:.1f} A{r},{r} 0 {large} 1 '
                           f'{x2:.1f},{y2:.1f} Z" fill="{colors[i % len(colors)]}"/>')
                mid = math.radians(ang + sweep / 2)
                lx, ly = cx + (r + 24) * math.cos(mid), cy + (r + 24) * math.sin(mid)
                pct = f"{100 * v / total:.0f}%"
                out.append(f'<text x="{lx:.0f}" y="{ly:.0f}" text-anchor="middle" font-size="12" '
                           f'fill="#111">{x(lb)} {pct}</text>')
                ang += sweep
        else:
            cw = (W - 2 * PAD) / len(values)
            base = H - PAD
            plot_h = H - 2 * PAD - 10
            for g in range(5):
                gy = base - plot_h * g / 4
                out.append(f'<line x1="{PAD}" y1="{gy:.0f}" x2="{W - PAD}" y2="{gy:.0f}" '
                           f'stroke="#e5e7eb"/>'
                           f'<text x="{PAD - 6}" y="{gy + 4:.0f}" text-anchor="end" font-size="10" '
                           f'fill="#666">{vmax * g / 4:g}</text>')
            if kind == "bar":
                for i, (lb, v) in enumerate(zip(labels, values)):
                    bh = plot_h * v / vmax
                    bx = PAD + i * cw + cw * 0.15
                    out.append(f'<rect x="{bx:.1f}" y="{base - bh:.1f}" width="{cw * 0.7:.1f}" '
                               f'height="{bh:.1f}" fill="{colors[i % len(colors)]}" rx="3"/>')
                    out.append(f'<text x="{PAD + i * cw + cw / 2:.1f}" y="{base - bh - 5:.0f}" '
                               f'text-anchor="middle" font-size="11" fill="#111">{v:g}</text>')
            else:  # line
                pts = []
                for i, v in enumerate(values):
                    px = PAD + i * cw + cw / 2
                    py = base - plot_h * v / vmax
                    pts.append(f"{px:.1f},{py:.1f}")
                out.append(f'<polyline points="{" ".join(pts)}" fill="none" stroke="#3b82f6" '
                           f'stroke-width="2.5"/>')
                for i, (p, v) in enumerate(zip(pts, values)):
                    px, py = p.split(",")
                    out.append(f'<circle cx="{px}" cy="{py}" r="4" fill="#3b82f6"/>'
                               f'<text x="{px}" y="{float(py) - 9:.0f}" text-anchor="middle" '
                               f'font-size="11" fill="#111">{v:g}</text>')
            for i, lb in enumerate(labels):
                out.append(f'<text x="{PAD + i * cw + cw / 2:.1f}" y="{H - PAD + 16}" '
                           f'text-anchor="middle" font-size="11" fill="#333">{x(lb)}</text>')
        out.append("</svg>")
        return "\n".join(out)

    def _exec_http(self, s: str) -> None:
        m = re.match(r"HTTP\s+([A-Za-z_][A-Za-z0-9_]*)\s*=\s*(GET|POST)\s+(.*)$", s, re.IGNORECASE | re.DOTALL)
        if not m:
            raise ScriptError(f"Bad HTTP (use: HTTP var = GET url | HTTP var = POST url BODY expr): {s}")
        var, method, rest = m.group(1), m.group(2).upper(), m.group(3)
        body = None
        bm = re.search(r"\s+BODY\s+(.+)$", rest, re.IGNORECASE | re.DOTALL)
        if bm:
            body = self.eval(bm.group(1))
            rest = rest[: bm.start()]
        url = str(self.eval(rest))
        if not url.startswith(("http://", "https://")):
            raise ScriptError(f"HTTP URL must start with http(s)://: {url}")
        self._log(f"HTTP {method} {url[:140]}")
        data = None
        headers = {"User-Agent": "SkillScript/1.0"}
        if method == "POST":
            data = (json.dumps(body) if isinstance(body, (dict, list)) else str(body or "")).encode()
            headers["Content-Type"] = "application/json"
        req = urllib.request.Request(url, data=data, headers=headers, method=method)
        try:
            with urllib.request.urlopen(req, timeout=30) as resp:
                self.vars[var] = resp.read(MAX_HTTP_BYTES).decode("utf-8", "replace")
                self.vars["HTTP_STATUS"] = resp.status
        except Exception as e:
            raise ScriptError(f"HTTP {method} failed: {e}")

    def _exec_file(self, s: str) -> None:
        u = s.upper()
        if u.startswith("FILE READ"):
            m = re.match(r"FILE\s+READ\s+([A-Za-z_][A-Za-z0-9_]*)\s*=\s*(.*)$", s, re.IGNORECASE)
            if not m:
                raise ScriptError(f"Bad FILE READ: {s}")
            p = self._safe_path(str(self.eval(m.group(2))))
            if not p.is_file():
                raise ScriptError(f"File not found: {p.name}")
            self.vars[m.group(1)] = p.read_text(encoding="utf-8", errors="replace")[:MAX_HTTP_BYTES]
            return
        mode = "w" if u.startswith("FILE WRITE") else ("a" if u.startswith("FILE APPEND") else None)
        if mode:
            rest = s[10 if mode == "w" else 11:]
            parts = self._split_top(rest, ",")
            if len(parts) < 2:
                raise ScriptError("FILE WRITE/APPEND needs: path, content")
            p = self._safe_path(str(self.eval(parts[0])))
            p.parent.mkdir(parents=True, exist_ok=True)
            with open(p, mode, encoding="utf-8") as f:
                f.write(str(self.eval(",".join(parts[1:]))))
            self._log(f"FILE {'WRITE' if mode == 'w' else 'APPEND'} {p.name}")
            return
        raise ScriptError(f"Unknown FILE operation: {s}")

    def _exec_mem(self, s: str) -> None:
        u = s.upper()
        if u.startswith("MEM SET"):
            parts = self._split_top(s[7:], ",")
            if len(parts) < 2:
                raise ScriptError("MEM SET needs: key, value")
            mem = _mem_load()
            mem[str(self.eval(parts[0]))] = self.eval(",".join(parts[1:]))
            _mem_save(mem)
            return
        if u.startswith("MEM GET"):
            m = re.match(r"MEM\s+GET\s+([A-Za-z_][A-Za-z0-9_]*)\s*=\s*(.*)$", s, re.IGNORECASE)
            if not m:
                raise ScriptError(f"Bad MEM GET: {s}")
            self.vars[m.group(1)] = _mem_load().get(str(self.eval(m.group(2))), "")
            return
        raise ScriptError(f"Unknown MEM operation: {s}")

    def _exec_email(self, s: str) -> None:
        m = re.match(r"EMAIL\s+TO\s+(.+?)\s+SUBJECT\s+(.+?)\s+BODY\s+(.+)$", s, re.IGNORECASE | re.DOTALL)
        if not m:
            raise ScriptError('Bad EMAIL (use: EMAIL TO expr SUBJECT expr BODY expr)')
        to = str(self.eval(m.group(1)))
        subject = str(self.eval(m.group(2)))
        body = str(self.eval(m.group(3)))
        if not self.email_sender:
            raise ScriptError("EMAIL unavailable in this context")
        self._log(f"EMAIL TO {to}: {subject[:80]}")
        self.email_sender(to, subject, body)


class _FlowSignal(Exception):
    pass


class _BreakSignal(_FlowSignal):
    pass


class _ContinueSignal(_FlowSignal):
    pass


class _ReturnSignal(_FlowSignal):
    pass


class _EndProgram(Exception):
    pass


# ------------------------------------------------------------- public API
def validate(code: str) -> dict:
    """Static validation: parse + block matching, no execution."""
    it = Interpreter(code)
    try:
        it._parse()
        if not it.program_name:
            return {"ok": False, "errors": ["Missing PROGRAM <name> statement on the first line"]}
        errors = []
        # verify every block terminator exists
        for i, (_, s) in enumerate(it.lines):
            u = s.upper()
            try:
                if u.startswith("IF ") and u.endswith("THEN"):
                    it._find_forward(i, "ENDIF", "END IF")
                elif u.startswith("FOR "):
                    it._find_forward(i, "NEXT")
                elif u.startswith("WHILE "):
                    it._find_forward(i, "WEND")
                elif u == "DO":
                    it._find_forward(i, "LOOP")
                elif u.startswith("SELECT CASE"):
                    it._find_forward(i, "END SELECT")
                elif u.startswith("SUB "):
                    it._find_forward(i, "END SUB")
                elif u.startswith("GOTO ") or u.startswith("GOSUB "):
                    lbl = s.split(None, 1)[1].strip().upper()
                    if lbl not in it.labels:
                        errors.append(f"Statement {i + 1}: unknown label {lbl}")
                elif u.startswith("CALL "):
                    if s[5:].strip().upper() not in it.subs:
                        errors.append(f"Statement {i + 1}: unknown SUB {s[5:].strip()}")
            except ScriptError as e:
                errors.append(f"Statement {i + 1}: {e}")
        return {"ok": not errors, "errors": errors, "name": it.program_name,
                "statements": len(it.lines), "labels": sorted(it.labels),
                "subs": sorted(it.subs)}
    except Exception as e:
        return {"ok": False, "errors": [str(e)]}


def run_script(code: str, context: Optional[dict] = None,
               allow_effects: bool = True) -> dict:
    """Execute a SkillScript program. Returns {ok, output, trace, ms} or
    {ok: False, error, trace}."""
    ai_runner = email_sender = None
    if allow_effects:
        def ai_runner(agent: str, prompt: str, system: str) -> str:  # noqa: F811
            from .providers import CodexProvider, ClaudeCodeProvider, ApiAgentProvider
            if agent == "claude":
                return ClaudeCodeProvider().run(prompt, system=system)
            if agent == "api":
                return ApiAgentProvider().run(prompt, system=system)
            return CodexProvider().run(prompt, system=system)

        def email_sender(to: str, subject: str, body: str) -> str:  # noqa: F811
            import uuid
            from .providers import SmtpEmailProvider
            from .config import get_config
            cfg = get_config()
            sender = cfg.get("smtp_from") or cfg.get("smtp_username") or "skillscript@local"
            r = SmtpEmailProvider().send(sender, [to], [], [], subject, body,
                                         idempotency_key=uuid.uuid4().hex)
            return getattr(r, "provider_message_id", "") or "sent"

    it = Interpreter(code, context=context, ai_runner=ai_runner, email_sender=email_sender)
    try:
        return it.run()
    except ScriptError as e:
        return {"ok": False, "error": str(e), "trace": it.trace,
                "output": "".join(it.output)}
    except Exception as e:
        return {"ok": False, "error": f"internal: {e}", "trace": it.trace,
                "output": "".join(it.output)}


def language_reference() -> str:
    """Human-readable language reference (served to the UI)."""
    return __doc__ or ""
